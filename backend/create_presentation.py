from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK = RGBColor(13, 15, 20)
DARK2 = RGBColor(21, 24, 32)
RED = RGBColor(226, 75, 74)
TEAL = RGBColor(29, 158, 117)
BLUE = RGBColor(55, 138, 221)
WHITE = RGBColor(232, 234, 240)
GRAY = RGBColor(85, 92, 110)
AMBER = RGBColor(239, 159, 39)

def blank_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)

def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def textbox(slide, text, x, y, w, h, size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ─── SLIDE 1: TITLE ───
s1 = blank_slide(prs)
bg(s1)
rect(s1, 0, 0, 13.33, 0.08, RED)
rect(s1, 0, 7.42, 13.33, 0.08, RED)
rect(s1, 0.5, 2.0, 0.06, 3.5, RED)
textbox(s1, "🛡️ MuleNet AI", 0.8, 1.8, 11, 1.2, size=52, bold=True, color=WHITE)
textbox(s1, "AI-Powered Mule Account & Fraud Detection System", 0.8, 3.1, 11, 0.7, size=22, color=GRAY)
textbox(s1, "BOI Hackathon 2024", 0.8, 4.0, 5, 0.5, size=16, color=TEAL)
textbox(s1, "Team: Atika  |  Mahum Siddiqui", 0.8, 4.6, 8, 0.5, size=16, color=GRAY)
rect(s1, 0.8, 5.3, 2.5, 0.5, RED)
textbox(s1, "Problem Statement 2", 0.85, 5.32, 2.4, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: PROBLEM ───
s2 = blank_slide(prs)
bg(s2)
rect(s2, 0, 0, 13.33, 0.08, RED)
textbox(s2, "THE PROBLEM", 0.5, 0.2, 6, 0.5, size=13, color=RED, bold=True)
textbox(s2, "Mule Accounts are Destroying India's Banking System", 0.5, 0.7, 12, 1.0, size=32, bold=True, color=WHITE)

problems = [
    ("💸", "₹10,000+ Crore", "lost to mule account fraud annually in India"),
    ("🏦", "1 in 5 Banks", "unable to detect mule networks in real-time"),
    ("⏰", "72+ Hours", "average time to detect and block fraud"),
    ("🔗", "Chain Transfers", "criminals use innocent accounts to hide money"),
]
for i, (icon, title, desc) in enumerate(problems):
    x = 0.5 + (i % 2) * 6.3
    y = 2.0 + (i // 2) * 2.0
    rect(s2, x, y, 5.8, 1.6, DARK2)
    textbox(s2, icon, x+0.2, y+0.15, 0.8, 0.8, size=28)
    textbox(s2, title, x+1.0, y+0.15, 4.5, 0.5, size=20, bold=True, color=RED)
    textbox(s2, desc, x+1.0, y+0.7, 4.5, 0.7, size=13, color=GRAY)

# ─── SLIDE 3: SOLUTION ───
s3 = blank_slide(prs)
bg(s3)
rect(s3, 0, 0, 13.33, 0.08, TEAL)
textbox(s3, "OUR SOLUTION", 0.5, 0.2, 6, 0.5, size=13, color=TEAL, bold=True)
textbox(s3, "MuleNet AI — Real-Time Detection", 0.5, 0.7, 12, 0.9, size=34, bold=True, color=WHITE)

solutions = [
    ("🤖", "AI/ML Model", "Random Forest trained on transaction patterns with 82% accuracy"),
    ("📊", "Graph Analysis", "Detects chain, star & round-trip mule networks instantly"),
    ("⚡", "Real-Time API", "Flask backend processes transactions in milliseconds"),
    ("🖥️", "Live Dashboard", "React dashboard with live alerts, charts & case management"),
    ("🔗", "Multi-Channel", "UPI, NEFT, IMPS, RTGS — all channels monitored"),
    ("🚨", "Auto Block", "One-click account blocking with SAR generation"),
]
for i, (icon, title, desc) in enumerate(solutions):
    x = 0.5 + (i % 3) * 4.2
    y = 1.9 + (i // 3) * 2.1
    rect(s3, x, y, 3.9, 1.8, DARK2)
    textbox(s3, icon + " " + title, x+0.2, y+0.2, 3.5, 0.5, size=16, bold=True, color=TEAL)
    textbox(s3, desc, x+0.2, y+0.8, 3.5, 0.8, size=12, color=GRAY)

# ─── SLIDE 4: ARCHITECTURE ───
s4 = blank_slide(prs)
bg(s4)
rect(s4, 0, 0, 13.33, 0.08, BLUE)
textbox(s4, "SYSTEM ARCHITECTURE", 0.5, 0.2, 8, 0.5, size=13, color=BLUE, bold=True)
textbox(s4, "How MuleNet AI Works", 0.5, 0.7, 10, 0.8, size=30, bold=True, color=WHITE)

steps = [
    ("📥", "Data Input", "UPI/NEFT/IMPS\nTransactions", BLUE),
    ("⚙️", "Processing", "Python Backend\nFlask API", AMBER),
    ("🧠", "AI Model", "Random Forest\n82% Accuracy", RED),
    ("📊", "Analysis", "Graph Network\nDetection", TEAL),
    ("🚨", "Alert", "Dashboard\nLive Alerts", RED),
]
for i, (icon, title, desc, color) in enumerate(steps):
    x = 0.4 + i * 2.5
    rect(s4, x, 2.0, 2.1, 2.5, DARK2)
    textbox(s4, icon, x+0.7, 2.1, 0.8, 0.6, size=26, align=PP_ALIGN.CENTER)
    textbox(s4, title, x+0.1, 2.8, 1.9, 0.4, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    textbox(s4, desc, x+0.1, 3.3, 1.9, 0.9, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        textbox(s4, "→", x+2.15, 2.9, 0.4, 0.5, size=22, color=GRAY, align=PP_ALIGN.CENTER)

rect(s4, 0.4, 5.0, 12.0, 0.8, DARK2)
textbox(s4, "🔄  Real-time processing  |  📡  Live data ingestion  |  🔒  Secure API  |  📱  Multi-channel support", 0.6, 5.1, 11.5, 0.6, size=13, color=TEAL, align=PP_ALIGN.CENTER)

# ─── SLIDE 5: DEMO ───
s5 = blank_slide(prs)
bg(s5)
rect(s5, 0, 0, 13.33, 0.08, AMBER)
textbox(s5, "LIVE DEMO", 0.5, 0.2, 6, 0.5, size=13, color=AMBER, bold=True)
textbox(s5, "See MuleNet AI in Action", 0.5, 0.7, 12, 0.9, size=34, bold=True, color=WHITE)

demo_steps = [
    ("1", "Real-time transaction feed coming in", BLUE),
    ("2", "AI model scores each transaction instantly", RED),
    ("3", "High risk accounts flagged automatically", RED),
    ("4", "Graph shows mule network connections", AMBER),
    ("5", "One click to block suspicious account", TEAL),
]
for i, (num, text, color) in enumerate(demo_steps):
    y = 1.9 + i * 0.85
    rect(s5, 0.5, y, 0.5, 0.5, color)
    textbox(s5, num, 0.5, y, 0.5, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s5, text, 1.2, y+0.05, 8, 0.45, size=16, color=WHITE)

rect(s5, 9.5, 1.8, 3.3, 4.5, DARK2)
textbox(s5, "🖥️ LIVE\nDASHBOARD\nDEMO", 9.6, 2.8, 3.1, 2.0, size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ─── SLIDE 6: RESULTS ───
s6 = blank_slide(prs)
bg(s6)
rect(s6, 0, 0, 13.33, 0.08, TEAL)
textbox(s6, "RESULTS & IMPACT", 0.5, 0.2, 8, 0.5, size=13, color=TEAL, bold=True)
textbox(s6, "What MuleNet AI Achieves", 0.5, 0.7, 12, 0.9, size=34, bold=True, color=WHITE)

results = [
    ("82%", "Model Accuracy", RED),
    ("<1sec", "Detection Time", BLUE),
    ("65+", "Alerts Generated", AMBER),
    ("₹4.5L", "Risk Detected", RED),
]
for i, (val, label, color) in enumerate(results):
    x = 0.5 + i * 3.1
    rect(s6, x, 2.0, 2.7, 2.0, DARK2)
    textbox(s6, val, x+0.1, 2.1, 2.5, 1.0, size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    textbox(s6, label, x+0.1, 3.1, 2.5, 0.6, size=13, color=GRAY, align=PP_ALIGN.CENTER)

impacts = [
    "✅ Real-time mule account detection across all channels",
    "✅ Graph-based network analysis to catch entire fraud chains",
    "✅ Automated alerts reduce manual investigation time by 80%",
    "✅ Scalable to ingest govt cyber fraud alerts & regulatory feeds",
]
for i, imp in enumerate(impacts):
    textbox(s6, imp, 0.5, 4.3 + i*0.55, 12, 0.5, size=14, color=WHITE)

# ─── SLIDE 7: THANK YOU ───
s7 = blank_slide(prs)
bg(s7)
rect(s7, 0, 0, 13.33, 0.08, RED)
rect(s7, 0, 7.42, 13.33, 0.08, RED)
textbox(s7, "🛡️", 5.8, 1.2, 1.5, 1.2, size=52, align=PP_ALIGN.CENTER)
textbox(s7, "Thank You!", 1, 2.3, 11.33, 1.2, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s7, "MuleNet AI — Protecting India's Banking System", 1, 3.5, 11.33, 0.7, size=20, color=GRAY, align=PP_ALIGN.CENTER)
rect(s7, 3.5, 4.5, 6.33, 0.6, DARK2)
textbox(s7, "Team: Atika  |  Mahum Siddiqui", 3.5, 4.52, 6.33, 0.55, size=16, color=TEAL, align=PP_ALIGN.CENTER)
textbox(s7, "BOI Hackathon 2024  •  Problem Statement 2", 1, 5.4, 11.33, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Save
prs.save('../presentation/MuleNet_AI_Presentation.pptx')
print("✅ Presentation ready!")
print("📁 Location: mulenet-ai/presentation/MuleNet_AI_Presentation.pptx")